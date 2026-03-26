from __future__ import annotations

import re
import requests
import json
import csv
import io

print("--- APP RELOADED: PDF FIXES APPLIED ---")
import os
import textwrap
import traceback
import zipfile
import time
import urllib.request
from pathlib import Path
from typing import Any, Dict, List, Tuple

from fastapi import FastAPI, Form, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import base64
import openai # For error catching
from reportlab.lib.pagesizes import LETTER
from reportlab.pdfgen import canvas
from reportlab.platypus import BaseDocTemplate, Frame, PageTemplate, Paragraph, Spacer, Image as PlatypusImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from pptx.util import Pt, Inches

# Optional AI integration (OpenAI, Azure OpenAI, Gemini, Bedrock)
# Global defaults (can be overridden by per-request config)
DEFAULT_AI_ENABLED = False
DEFAULT_AI_PROVIDER = os.getenv("AI_PROVIDER", "auto").lower()
DEFAULT_AI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
# ... other env vars ...

AI_STATUS_MSG = "AI not initialized"
AI_DEBUG = True 

# We'll use a Factory or local initialization for dynamic requests
try:
    from openai import AzureOpenAI, OpenAI
except ImportError:
    AzureOpenAI = None
    OpenAI = None

try:
    import google.generativeai as genai
except ImportError:
    genai = None

try:
    import boto3
except ImportError:
    boto3 = None


def _get_start_client(provider: str, key: str, model: str, endpoint: str = None) -> Any:
    """Initialize an AI client dynamically."""
    if provider == "openai":
        if not OpenAI: return None
        return OpenAI(api_key=key)
    elif provider == "openrouter":
        if not OpenAI: return None
        return OpenAI(
            api_key=key,
            base_url="https://openrouter.ai/api/v1",
            default_headers={
                "HTTP-Referer": "https://asthra-app.com", # Required by OpenRouter
                "X-Title": "ASTHRA",
            }
        )
    elif provider == "gemini":
        if not genai: return None
        genai.configure(api_key=key)
        return genai
    elif provider == "bedrock":
        if not boto3: return None
        return boto3.client("bedrock-runtime", region_name=os.getenv("AWS_DEFAULT_REGION", "us-east-1"))
    return None

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

BASE_DIR = Path(__file__).resolve().parent
FILES_DIR = BASE_DIR / "generated_files"
FILES_DIR.mkdir(exist_ok=True)


class RegenerateRequest(BaseModel):
    plan: Dict[str, Any]
    ppt_settings: Dict[str, Any] | None = None
    cert_settings: Dict[str, Any] | None = None

class AiAssistRequest(BaseModel):
    text: str
    instruction: str # "Summarize", "Expand", "Formalize"
    api_key: str | None = None
    provider: str | None = None
    model: str | None = None


def _fetch_openrouter_models(only_free=False) -> List[Dict[str, Any]]:
    try:
        url = "https://openrouter.ai/api/v1/models"
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req) as response:
            if response.status == 200:
                data = json.loads(response.read().decode())
                models = []
                for m in data.get("data", []):
                    pricing = m.get("pricing", {})
                    is_free = (
                        pricing.get("prompt") == "0" and 
                        pricing.get("completion") == "0"
                    ) or "free" in m["id"]
                    
                    if only_free and not is_free:
                        continue
                        
                    models.append({
                        "id": m["id"],
                        "name": m["name"],
                        "is_free": is_free
                    })
                PRIORITY_MODELS = [
                    "google/gemini-2.0-flash-exp:free",
                    "google/gemini-2.0-flash-thinking-exp:free",
                    "meta-llama/llama-3.3-70b-instruct:free",
                    "google/gemini-exp-1206:free",
                    "mistralai/mistral-nemo:free",
                ]
                
                def sort_key(x):
                    # 1. Priority (Lower index = Better)
                    try:
                        p_index = PRIORITY_MODELS.index(x["id"])
                    except ValueError:
                        p_index = 999
                    
                    # 2. Free vs Paid (True=Free first) -> False < True, so we want not is_free for ascending?
                    # We want Free first? False (0) comes before True (1).
                    # Actually: not True = False(0). not False = True(1).
                    # So (0, ...) comes first.
                    
                    return (p_index, not x["is_free"], x["name"])

                models.sort(key=sort_key)
                return models
    except:
        pass
    return []

@app.get("/models")
async def list_models(provider: str = "openrouter"):
    """Fetch available models. For OpenRouter, we fetch from their API."""
    if provider == "openrouter":
        models = _fetch_openrouter_models()
        return {"models": models}
    
    return {"models": []}


def _fallback_plan(message: str, file: UploadFile | None = None) -> Dict[str, Any]:
    file_note = f"Uploaded file: {file.filename}" if file else "No upload provided."
    return {
        "title": "ASTHRA Project Documentation",
        "summary": f"Generated documentation for: {message}",
        "sections": [
            {
                "heading": "Overview",
                "bullets": [
                    "Project statement received and logged.",
                    "ASTHRA generated offline demo content.",
                    file_note,
                ],
            },
        ],
        "claims": ["Automated generation demo."],
        "certificate_note": "Demo certificate.",
    }

def _greeting_plan(message: str) -> Dict[str, Any]:
    return {
        "title": "Welcome to ASTHRA",
        "summary": "Hi there! I can help you generate documentation.",
        "sections": [{"heading": "Help", "bullets": ["Enter a prompt to start."]}],
        "claims": ["Interactive assistant."],
        "certificate_note": "Awaiting project details.",
    }

def _is_greeting(message: str) -> bool:
    msg = message.strip().lower()
    msg = msg.strip(" .,!?:;")
    # Basic greetings
    if msg in {"hi", "hey", "hello", "hiya", "yo", "sup", "hola", "greetings"}:
        return True
    # Starts with greeting
    if msg.startswith(("hi ", "hello ", "hey ")):
        return True
    # Conversational questions that shouldn't trigger project generation
    common_questions = {
        "how are you", "how are you doing", "what is up", "whats up", 
        "who are you", "what are you", "what is your name"
    }
    if any(q in msg for q in common_questions):
        return True
        
    return False


def _extract_github_url(text: str) -> str | None:
    # Basic regex to find github.com/user/repo
    pattern = r"https?://github\.com/[a-zA-Z0-9_\-]+/[a-zA-Z0-9_\-]+"
    match = re.search(pattern, text)
    if match:
        return match.group(0)
    return None

def _fetch_github_context(url: str) -> str:
    """Fetch README and structure from public GitHub repo, plus source code for analysis."""
    try:
        # Convert HTML URL to API URL
        api_url = url.replace("github.com/", "api.github.com/repos/")
        if api_url.endswith(".git"):
            api_url = api_url[:-4]

        headers = {"Accept": "application/vnd.github.v3+json", "User-Agent": "Asthra-Bot"}
        
        # 1. Fetch Repo Info
        info_resp = requests.get(api_url, headers=headers)
        if info_resp.status_code != 200:
            return f"[Error fetching GitHub repo: {info_resp.status_code}]"

        data = info_resp.json()
        description = data.get("description", "No description.")
        
        # 2. Fetch README
        readme_url = f"{api_url}/readme"
        readme_resp = requests.get(readme_url, headers=headers)
        readme_content = ""
        if readme_resp.status_code == 200:
            try:
                readme_data = readme_resp.json()
                readme_text = base64.b64decode(readme_data["content"]).decode("utf-8")
                readme_content = readme_text[:3000] + ("\n...(truncated)" if len(readme_text) > 3000 else "")
            except:
                readme_content = "(Could not decode README)"
        
        # 3. Fetch Code Files (Depth exploration)
        code_context = []
        total_chars = 0
        MAX_CHARS = 25000 # Limit to avoid context overflow
        ALLOWED_EXTS = {".py", ".js", ".ts", ".dart", ".java", ".c", ".cpp", ".h", ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt"}
        
        def fetch_dir(url, depth=0):
            nonlocal total_chars
            if depth > 2 or total_chars >= MAX_CHARS: return []
            
            resp = requests.get(url, headers=headers)
            if resp.status_code != 200: return []
            
            items = resp.json()
            # Prioritize certain folders
            items.sort(key=lambda x: 0 if x['name'] in ('src', 'lib', 'app', 'backend') else 1)
            
            collected = []
            
            # First pass: Files
            for item in items:
                if total_chars >= MAX_CHARS: break
                
                if item['type'] == 'file':
                    ext = os.path.splitext(item['name'])[1]
                    if ext in ALLOWED_EXTS:
                        # Fetch Content
                        file_resp = requests.get(item['url'], headers=headers)
                        if file_resp.status_code == 200:
                            try:
                                f_data = file_resp.json()
                                content = base64.b64decode(f_data['content']).decode('utf-8')
                                if len(content) > 6000:
                                    content = content[:6000] + "\n...(file truncated)"
                                
                                formatted = f"\n--- FILE: {item['path']} ---\n{content}\n"
                                collected.append(formatted)
                                total_chars += len(formatted)
                            except: pass
            
            # Second pass: Directories (only if space remains)
            for item in items:
                if total_chars >= MAX_CHARS: break
                if item['type'] == 'dir' and item['name'] not in {'.git', 'node_modules', 'venv', 'test', 'tests'}:
                    collected.extend(fetch_dir(item['url'], depth + 1))
            
            return collected

        # Start fetching from contents URL
        if total_chars < MAX_CHARS:
            contents_url = f"{api_url}/contents"
            code_context.extend(fetch_dir(contents_url))

        full_code = "".join(code_context)
                 
        context = f"""
        GITHUB REPOSITORY CONTEXT:
        URL: {url}
        Description: {description}
        
        README (Summary):
        {readme_content}
        
        SOURCE CODE ANALYSIS DATA:
        {full_code}
        ------------------------------------------
        """
        return context

    except Exception as e:
        print(f"GitHub Fetch Error: {e}")
        return f"[Failed to read GitHub repo: {e}]"



def _system_prompt() -> str:
    return """You are an expert technical consultant and documentation specialist. Your goal is to draft a comprehensive, professional, and detailed project report. The output must be a valid JSON object.

CRITICAL INSTRUCTION: The 'sections' array MUST contain OBJECTS, not strings. Each object in 'sections' MUST have a 'heading' and a 'bullets' list. Do NOT output sections as a list of strings.

Structure:
{
  "title": "Project Title",
  "summary": "A detailed executive summary (approx 100-150 words).",
  "sections": [
    {
      "heading": "Section Title (e.g., Introduction, Methodology, System Design)",
      "bullets": ["Detailed paragraph 1...", "Detailed paragraph 2..."]
    }
  ],
  "ppt_slides": [
    { "title": "Title Slide", "bullets": ["Key point 1", "Key point 2"] }
  ],
  "claims": ["Patent-style claim 1...", "Patent-style claim 2..."],
  "certificate_note": "A professional citation for the certificate."
}

REQUIREMENTS:
1. Write in a formal, academic tone.
2. 'sections' are for the REPORT (Detailed, paragraphs).
3. 'ppt_slides' are for the PRESENTATION (Concise, bullet points, fewer words).
4. Ensure 'sections' is a list of objects with populated 'bullets'.
5. You MUST provide at least 2 content bullets per section. Empty bullets are invalid.
6. Respond with JSON ONLY.
"""


def _call_ai_plan_dynamic(
    message: str, 
    provider: str, 
    api_key: str, 
    model: str
) -> Tuple[Dict[str, Any] | None, str | None]:
    """Call AI with dynamic credentials."""
    
    client_instance = _get_start_client(provider, api_key, model)
    if not client_instance:
        return None, f"Could not initialize provider {provider}"

    try:
        # OPENAI / OPENROUTER
        if provider in {"openai", "openrouter"}:
            
            # For OpenRouter, we implement a fallback loop for Rate Limits
            candidate_models = [model]
            if provider == "openrouter":
                # If it's openrouter, fetch other free models to try if the first one fails
                try:
                    free_models = _fetch_openrouter_models(only_free=True)
                    # Add them to candidates, avoiding duplicates
                    for fm in free_models:
                        if fm["id"] != model:
                            candidate_models.append(fm["id"])
                except:
                    pass

            last_error = None
            
            for attempt_model in candidate_models:
                try:
                    if AI_DEBUG: print(f"Using {provider} model={attempt_model}")
                    # Re-init client if needed or just use the same client with different model param
                    
                    response = client_instance.chat.completions.create(
                        model=attempt_model,
                        temperature=0.5,
                        max_tokens=4000,
                        response_format={"type": "json_object"},
                        messages=[
                            {"role": "system", "content": _system_prompt()},
                            {"role": "user", "content": message},
                        ],
                    )
                    raw_content = response.choices[0].message.content
                    return _parse_json_content(raw_content)
                    
                except openai.RateLimitError as e:
                    print(f"RateLimit hit for {attempt_model}: {e}. Switching...")
                    last_error = e
                    continue # Try next model
                except Exception as e:
                    # For other errors, we might not want to switch, or maybe we do?
                    # Let's assume other errors are fatal or model-specific, so we can try next.
                    print(f"Error with {attempt_model}: {e}")
                    last_error = e
                    continue
            
            return None, f"All models failed. Last error: {str(last_error)}"

        # GEMINI
        if provider == "gemini":
            model_instance = client_instance.GenerativeModel(model)
            response = model_instance.generate_content(
                f"{_system_prompt()}\nUser request:\n{message}\nReply with JSON only.",
                 generation_config=genai.types.GenerationConfig(
                    max_output_tokens=4000
                ) if genai else None
            )
            return _parse_json_content(response.text)

        # BEDROCK (Keeping basic logic)
        if provider == "bedrock":
            prompt = {
                "messages": [{"role": "user", "content": [{"type": "text", "text": f"{_system_prompt()}\nUser request:\n{message}"}]}],
                "max_tokens": 4000,
                "temperature": 0.5,
                "anthropic_version": "bedrock-2023-05-31",
            }
            result = client_instance.invoke_model(modelId=model, body=json.dumps(prompt))
            payload = json.loads(result["body"].read())
            raw_content = payload["content"][0]["text"]
            return _parse_json_content(raw_content)

    except Exception as exc:
        if AI_DEBUG: traceback.print_exc()
        return None, str(exc)

    return None, "Provider logic fallthrough"


def _normalize_plan(plan: Dict[str, Any], fallback_message: str) -> Dict[str, Any]:
    title = plan.get("title") or "ASTHRA Project Documentation"
    summary = plan.get("summary") or f"Generated documentation for: {fallback_message}"
    sections = []
    for section in plan.get("sections", []):
        if isinstance(section, dict):
            heading = section.get("heading") or "Section"
            bullets = [b.strip() for b in (section.get("bullets") or []) if str(b).strip()]
        else:
            # Fallback if section is malformed (e.g. just a string)
            heading = str(section)
            bullets = ["Details pending."]
        sections.append({"heading": heading, "bullets": bullets or ["Details pending."]})

    if not sections:
        sections = [{"heading": "Overview", "bullets": [summary]}]

    claims = [c.strip() for c in plan.get("claims", []) if str(c).strip()] or ["Automated documentation generation."]
    certificate_note = plan.get("certificate_note") or "Generated via ASTHRA."

    return {
        "title": title,
        "summary": summary,
        "sections": sections,
        "claims": claims,
        "certificate_note": certificate_note,
    }


def _parse_json_content(raw_content: str) -> Tuple[Dict[str, Any] | None, str | None]:
    if not raw_content: return None, "Empty AI response"
    cleaned = raw_content.strip()
    for prefix in ("```json", "```"):
        if cleaned.lower().startswith(prefix):
            cleaned = cleaned[len(prefix) :].strip()
    if cleaned.endswith("```"):
        cleaned = cleaned[:-3].strip()
    try:
        data = json.loads(cleaned)
        # Fallback for ppt_slides if missing (old plans or smaller models)
        if "ppt_slides" not in data or not data["ppt_slides"]:
             data["ppt_slides"] = data.get("sections", [])
        return data, None
    except json.JSONDecodeError:
        return None, f"Non-JSON AI response: {cleaned[:200]}"


def _wrap_lines(text: str, width: int) -> List[str]:
    return textwrap.wrap(text, width=width) if text else []


def _build_report_pdf(plan: Dict[str, Any], path: Path) -> None:
    print(f"DEBUG: Building PDF w/ sections: {len(plan.get('sections', []))} - Title: {plan.get('title')}")
    doc = BaseDocTemplate(str(path), pagesize=LETTER)
    
    # IEEE Styles
    styles = getSampleStyleSheet()
    
    # Title Style: Centered, Large, Serif
    style_title = ParagraphStyle(
        'IEEE_Title', 
        parent=styles['Heading1'], 
        fontName='Times-Bold', 
        fontSize=24, 
        alignment=TA_CENTER, 
        spaceAfter=20
    )
    
    # Abstract/Summary: Centered, Italic, Serif
    style_abstract = ParagraphStyle(
        'IEEE_Abstract',
        parent=styles['Normal'],
        fontName='Times-Italic',
        fontSize=10,
        alignment=TA_JUSTIFY,
        leftIndent=40,
        rightIndent=40,
        spaceAfter=20
    )
    
    # Headings: Bold, Serif
    style_heading = ParagraphStyle(
        'IEEE_Heading', 
        parent=styles['Heading2'], 
        fontName='Times-Bold', 
        fontSize=12, 
        spaceBefore=10, 
        spaceAfter=4
    )
    
    # Body: Serif, Justified, 2-Column friendly
    style_body = ParagraphStyle(
        'IEEE_Body', 
        parent=styles['Normal'], 
        fontName='Times-Roman', 
        fontSize=10, 
        alignment=TA_JUSTIFY,
        leading=12
    )
    
    # Define Frames
    # 1. Title Page Frame (Full width for Title/Abstract) - actually we can just use normal flow 
    # but strictly speaking 2-col usually starts after title.
    # Let's define a "TwoColumn" Template.
    
    frame_left = Frame(0.5*inch, 0.5*inch, 3.5*inch, 10*inch, id='col1')
    frame_right = Frame(4.25*inch, 0.5*inch, 3.5*inch, 10*inch, id='col2')
    
    page_template = PageTemplate(id='TwoCol', frames=[frame_left, frame_right])
    doc.addPageTemplates([page_template])
    
    story = []
    
    # Title & Abstract (Spanning attempts? ReportLab Frames flow into next frame. 
    # To span columns, we need a separate "Title" template or just accept flow)
    # SIMPLE APPROACH for mixed layout:
    # We will use a custom FirstPageTemplate if we want Title to span.
    # BUT easier: Just define Title as "Flowable" that is wide? No, frames clip.
    # Correct Way: Define a Single Column Frame for the top, then Two Columns below? 
    # Or just use single col for first page? 
    # IEEE standard: Title spans, text is 2-col.
    
    # Let's switch to a custom FirstPageTemplate
    frame_full = Frame(0.5*inch, 0.5*inch, 7.5*inch, 10*inch, id='full')
    template_cover = PageTemplate(id='Cover', frames=[frame_full])
    
    # We need to register both, verify switching.
    # Actually, simpler hack for "AI" generation:
    # Just render title in the first column? No, looks bad.
    # Let's stick to: Entire document is 2-column, Title is just big text in left column? No.
    # We will use 'Cover' for first page? But text needs to flow immediately to 2-col?
    # That requires "NextPageTemplate".
    # Implementation:
    # 1. Add 'Cover' template and 'TwoCol' template.
    # 2. Start with 'Cover'.
    # 3. Add Title, Summary.
    # 4. Add "NextPageTemplate('TwoCol')" (Will apply to NEXT page? Or force break?)
    # 5. Ideally, we want Title then immediate columns. This requires a complex Frame setup on Page 1.
    #    (FrameTop, FrameBottomLeft, FrameBottomRight).
    
    # FRAMES SETUP
    # Page 1: Title (Top 2 inches), and 2 Columns (Bottom 8 inches)
    f_title = Frame(0.5*inch, 8.5*inch, 7.5*inch, 2*inch, id='title_frame')
    f_col1 = Frame(0.5*inch, 0.5*inch, 3.6*inch, 7.8*inch, id='col1_p1')
    f_col2 = Frame(4.4*inch, 0.5*inch, 3.6*inch, 7.8*inch, id='col2_p1')
    
    template_p1 = PageTemplate(id='FirstPage', frames=[f_title, f_col1, f_col2])
    
    # Page 2+: Full height 2 cols
    f_left = Frame(0.5*inch, 0.5*inch, 3.6*inch, 10*inch, id='left')
    f_right = Frame(4.4*inch, 0.5*inch, 3.6*inch, 10*inch, id='right')
    template_normal = PageTemplate(id='Normal', frames=[f_left, f_right])
    
    doc.addPageTemplates([template_p1, template_normal])
    
    # Content
    story.append(Paragraph(plan["title"], style_title))
    story.append(Paragraph(f"<b>Abstract:</b> {plan['summary']}", style_abstract))
    
    # Force jump to next frame (which is col1_p1)? 
    # ReportLab fills frames in order. f_title is first.
    # So once title/abstract are done, it naturally flows to f_col1_p1? 
    # Yes, provided they fit. 
    # If title is huge, it might overflow. We assume it fits in 2 inches.
    
    story.append(Spacer(1, 0.2*inch))
    
    for section in plan["sections"]:
        story.append(Paragraph(section["heading"], style_heading))
        for bullet in section["bullets"]:
            clean_b = bullet.replace('\n', ' ').strip()
            story.append(Paragraph(f"• {clean_b}", style_body))
        story.append(Spacer(1, 0.1*inch))
        
    doc.build(story)




def _build_docx(plan: Dict[str, Any], path: Path, template_path: Path | None = None) -> None:
    if template_path and template_path.exists():
        try:
            doc = Document(template_path)
            # Placeholder Replacement Strategy
            
            # Construct content
            full_content = ""
            for sec in plan.get("sections", []):
                full_content += f"{sec['heading']}\n"
                for b in sec.get("bullets", []):
                    full_content += f"- {b}\n"
                full_content += "\n"
            
            replacements = {
                "{{title}}": plan.get("title", ""),
                "{{summary}}": plan.get("summary", ""),
                "{{abstract}}": plan.get("summary", ""),
                "{{content}}": full_content,
                "{{body}}": full_content,
                "{{date}}": time.strftime("%Y-%m-%d"),
            }
            
            # Replace in Paragraphs
            for p in doc.paragraphs:
                for k, v in replacements.items():
                    if k in p.text:
                        p.text = p.text.replace(k, v)
            
            # Replace in Tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                             for k, v in replacements.items():
                                if k in p.text:
                                    p.text = p.text.replace(k, v)
                                    
            doc.save(str(path))
            return
        except Exception as e:
            print(f"Template Docx Error: {e}")
            # Fallback to default
            
    doc = Document()
    doc.add_heading(plan["title"], 0)
    
    doc.add_heading('Executive Summary', level=1)
    doc.add_paragraph(plan["summary"])
    
    for section in plan["sections"]:
        doc.add_heading(section["heading"], level=1)
        for bullet in section["bullets"]:
            doc.add_paragraph(bullet, style='List Bullet')
            
    doc.add_heading('Patent Claims', level=1)
    for idx, claim in enumerate(plan["claims"], 1):
        doc.add_paragraph(f"{idx}. {claim}", style='List Number')
        
    doc.save(str(path))



def _draw_header(c, text):
    c.saveState()
    c.setFont("Helvetica-Bold", 16)
    c.drawString(40, 750, text)
    c.line(40, 745, 570, 745)
    c.restoreState()
    return 720

def _draw_paragraph(c, text, y):
    # Simple wrapping
    text_obj = c.beginText(40, y)
    text_obj.setFont("Helvetica", 11)
    
    words = text.split(' ')
    line = ""
    for word in words:
        if c.stringWidth(line + word, "Helvetica", 11) < 500:
            line += word + " "
        else:
            text_obj.textLine(line)
            line = word + " "
            y -= 14
    text_obj.textLine(line)
    y -= 14
    
    c.drawText(text_obj)
    return y - 20

def _build_patent_pdf(plan: Dict[str, Any], message: str, path: Path) -> None:
    c = canvas.Canvas(str(path), pagesize=LETTER)
    y = _draw_header(c, "ASTHRA Patent Draft")
    c.setFont("Helvetica", 11)
    y = _draw_paragraph(c, f"Based on: {message}", y)
    c.setFont("Helvetica-Bold", 13)
    c.drawString(60, y, "Claims")
    y -= 18
    c.setFont("Helvetica", 11)
    for idx, claim in enumerate(plan["claims"], start=1):
        for line in _wrap_lines(f"{idx}. {claim}", 90):
            c.drawString(70, y, line)
            y -= 14
    c.save()


def _build_certificates_zip(plan: Dict[str, Any], message: str, path: Path, names: List[str] = [], data: List[Dict[str, str]] = None, settings: Dict[str, Any] | None = None) -> None:
    # Normalize Data
    settings = settings or {}
    rows = []
    if data:
        rows = data
    elif names:
        rows = [{"name": n} for n in names]
    else:
        rows = [{"name": f"Participant {i}"} for i in range(1, 4)]
        
    # Elements Engine
    elements = settings.get("elements", [])
    bg_image_b64 = settings.get("background_image") if settings else None
    
    # If no elements provided, create default set (Backward Compatibility)
    # If no elements provided, create default set (Backward Compatibility)
    if not elements:
        # CHECK IF SETTINGS HAS ELEMENTS (From Designer)
        if settings and "elements" in settings:
            elements = settings["elements"]
        else:
            # Defaults
            title = settings.get("cert_title", "Certificate of Completion") if settings else "Certificate of Completion"
            event = settings.get("event_name", plan.get("title", "Project")) if settings else plan.get("title", "Project")
            auth = settings.get("authority_name", "Program Director") if settings else "Program Director"
            color = settings.get("theme_color", "#000000") if settings else "#000000"
            
            # We use old layout logic to map to new elements
            layout = settings.get("layout", {})
            def gc(k, d): 
                 try: return int(layout.get(k, d))
                 except: return d
            
            elements = [
                {"type": "text", "text": title, "x": gc("title_x", 421), "y": gc("title_y", 480), "size": 36, "font": "Helvetica-Bold", "color": color, "align": "center"},
                {"type": "text", "text": "This is presented to", "x": gc("presented_x", 421), "y": gc("presented_y", 420), "size": 18, "font": "Helvetica", "color": "#666666", "align": "center"},
                {"type": "text", "text": "{name}", "x": gc("name_x", 421), "y": gc("name_y", 350), "size": 40, "font": "Helvetica-Bold", "color": "#000000", "align": "center"},
                {"type": "text", "text": "For successful contribution to:", "x": gc("mk_x", 421), "y": gc("mk_y", 280), "size": 18, "font": "Helvetica", "color": "#666666", "align": "center"},
                {"type": "text", "text": event, "x": gc("event_x", 421), "y": gc("event_y", 230), "size": 24, "font": "Helvetica-Bold", "color": color, "align": "center"},
                {"type": "text", "text": f"Date: {message}", "x": gc("date_x", 100), "y": gc("date_y", 100), "size": 14, "font": "Helvetica", "color": "#000000", "align": "left"},
                {"type": "text", "text": auth, "x": gc("auth_x", 650), "y": gc("auth_y", 95), "size": 16, "font": "Helvetica", "color": "#000000", "align": "center"}
            ]

    # Helper to draw
    def draw_elements(c, data_context):
        # Draw Background 
        if bg_image_b64:
             try:
                 bg_bytes = base64.b64decode(bg_image_b64)
                 with open("tmp_bg.png", "wb") as f: f.write(bg_bytes)
                 c.drawImage("tmp_bg.png", 0, 0, width=842, height=595)
             except: pass
        else:
             # Default Border if no bg
             try:
                # get color from first element or default
                c_hex = elements[0].get("color", "#000000")
                use_rgb = hex_to_rgb_tuple(c_hex)
                c.setStrokeColorRGB(*use_rgb)
                c.setLineWidth(5)
                c.rect(20, 20, 802, 555)
             except: pass

        for el in elements:
            etype = el.get("type", "text")
            x = int(el.get("x", 0))
            y = int(el.get("y", 0))
            
            if etype == "text":
                text = el.get("text", "")
                # Replace placeholders
                for k, v in data_context.items():
                    text = text.replace(f"{{{k}}}", str(v))
                
                font = el.get("font", "Helvetica")
                size = int(el.get("size", 12))
                color_hex = el.get("color", "#000000")
                
                rgb = hex_to_rgb_tuple(color_hex)
                c.setFillColorRGB(*rgb)
                c.setFont(font, size)
                
                align = el.get("align", "left")
                if align == "center":
                    c.drawCentredString(x, y, text)
                elif align == "right":
                    c.drawRightString(x, y, text)
                else:
                    c.drawString(x, y, text)
                    
            elif etype == "image":
                img_b64 = el.get("value")
                w = int(el.get("width", 100))
                h = int(el.get("height", 100))
                if img_b64:
                    try:
                        img_bytes = base64.b64decode(img_b64)
                        fname = f"tmp_el_{time.time()}.png"
                        with open(fname, "wb") as f: f.write(img_bytes)
                        c.drawImage(fname, x, y, width=w, height=h, mask='auto')
                        try: os.remove(fname)
                        except: pass
                    except: pass
    
    # helper for hex
    def hex_to_rgb_tuple(h):
        try:
            h = h.lstrip('#')
            return tuple(int(h[i:i+2], 16)/255.0 for i in (0, 2, 4))
        except: return (0,0,0)

    # Generate Loop
    with zipfile.ZipFile(path, 'w') as zipf:
        for row in rows:
            # Determine filename from 'name' or first value
            p_name = row.get("name", list(row.values())[0] if row else "Participant")
            safe_name = "".join([c for c in str(p_name) if c.isalnum() or c in (' ', '_')]).strip()
            cert_filename = f"Certificate_{safe_name}.pdf"
            cert_path = FILES_DIR / cert_filename
            
            c = canvas.Canvas(str(cert_path), pagesize=(842, 595))
            
            # Context - Merge row with global defaults
            ctx = {"date": message, "title": plan.get("title", "")}
            for k, v in row.items():
                ctx[k] = v
            
            draw_elements(c, ctx)
            c.save()
            zipf.write(cert_path, cert_filename)
    
    try: os.remove("tmp_bg.png")
    except: pass
def _build_ppt(plan: Dict[str, Any], path: Path, custom_template_path: Path | None = None, settings: Dict[str, Any] | None = None, template_path: Path | None = None) -> None:
    # Logic: 
    # 1. Custom template uploaded (template_path)? Use it and replace placeholders.
    # 2. Existing 'template_modern.pptx'? Use it.
    # 3. Default white.
    
    # Unified template path check
    use_template = template_path or custom_template_path
    
    if use_template and use_template.exists():
        try:
            prs = Presentation(str(use_template))
            # Placeholder Strategy: {{title}}, {{date}} in existing slides
            replacements = {
                "{{title}}": plan.get("title", ""),
                "{{summary}}": plan.get("summary", ""),
                "{{date}}": time.strftime("%Y-%m-%d"),
            }
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            if "{{" in p.text:
                                for k, v in replacements.items():
                                    if k in p.text:
                                        p.text = p.text.replace(k, v)
                                        
            # Add New Slides Strategy:
            # We want to add new content slides using the template's layout.
            # Usually Layout 1 is "Title and Content".
            bullet_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            
            slides_source = plan.get("ppt_slides", plan.get("sections", []))
            for section in slides_source:
                 slide = prs.slides.add_slide(bullet_layout)
                 try:
                     # Attempt to set Title
                     if slide.shapes.title:
                        slide.shapes.title.text = section.get("heading", section.get("title", ""))
                 except: pass
                 
                 # Attempt to find body placeholder
                 body = None
                 for ph in slide.placeholders:
                     if ph.placeholder_format.idx == 1: # Content
                         body = ph
                         break
                 
                 if body and body.has_text_frame:
                     tf = body.text_frame
                     tf.text = ""
                     for b in section["bullets"]:
                         p = tf.add_paragraph()
                         p.text = b
                         p.level = 0
            
            prs.save(str(path))
            return

        except Exception as e:
            print(f"Failed to load custom template: {e}")
            
    # Default Logic Fallback
    prs = None
    fallback_tmpl = BASE_DIR / "templates" / "template_modern.pptx"
    if fallback_tmpl.exists():
        try: prs = Presentation(str(fallback_tmpl))
        except: prs = Presentation()
    else:
        prs = Presentation()

    # Apply Settings (Font & Color & Background)
    # IEEE Standard: Times New Roman or Arial.
    target_font = settings.get("font_name", "Times New Roman") if settings else "Times New Roman"
    title_color_hex = settings.get("title_color", "000000") if settings else "000000"
    body_color_hex = settings.get("body_color", "000000") if settings else "000000"
    bg_image_b64 = settings.get("background_image") if settings else None

    # Handle Background Image
    bg_image_path = None
    if bg_image_b64:
        try:
            # Decode base64 to temp file
            bg_data = base64.b64decode(bg_image_b64)
            # Detect extension? Hard with base64 without header. 
            # We assume it's image data readable by ReportLab (Auto-detects)
            # But we need a file path. Let's use a generic name but maybe .img or .png works?
            # ReportLab checks magic numbers usually.
            # But python-docx/pptx might be pickier.
            # Let's try to detect magic number or just default to .png (most work as png or jpg)
            # Better: Write as temporary file with no extension/generic, reportlab handles it.
            # PPTX needs known extension.
            # We will try to guess from magic bytes or just use .png which often works for mixed types in some libs,
            # BUT better is to check first bytes.
            
            ext = ".png" # Default
            if bg_data.startswith(b'\xff\xd8'): ext = ".jpg"
            elif bg_data.startswith(b'\x89PNG'): ext = ".png"
            elif bg_data.startswith(b'BMP'): ext = ".bmp"
            elif bg_data.startswith(b'RIFF') and b'WEBP' in bg_data[0:20]: ext = ".webp"
            
            bg_image_path = FILES_DIR / f"temp_bg_image{ext}"
            with open(bg_image_path, "wb") as f:
                f.write(bg_data)
        except Exception as e:
            print(f"Failed to process background image: {e}")

    def _hex_to_rgb(hex_str):
        if not hex_str: return None
        hex_str = hex_str.lstrip('#')
        try:
            return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
        except: return RGBColor(0, 0, 0)

    def _apply_formatting(shape, is_title=False):
        if not shape.has_text_frame:
            return
            
        font_name = target_font
        color_hex = title_color_hex if is_title else body_color_hex
        rgb = _hex_to_rgb(color_hex)
        
        for paragraph in shape.text_frame.paragraphs:
             # Just set paragraph level properties if possible, or run level
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.color.rgb = rgb
            # Set font for new runs too? 
            # We can set default properties on text_frame sometimes, but runs override.
            
    def _add_slide_number(slide, number):
        # Add footer with slide number
        txBox = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(1), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(12)
        p.font.name = target_font
        p.font.color.rgb = RGBColor(100, 100, 100)

    # TITLES
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # 1. Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = plan["title"]
    subtitle = slide.placeholders[1]
    subtitle.text = plan["summary"]
    _apply_formatting(slide.shapes.title, is_title=True)
    _apply_formatting(subtitle, is_title=False)
    
    # USE SEPARATE PPT SLIDES if available
    slides_source = plan.get("ppt_slides", plan.get("sections", []))
    
    slide_idx = 1
    for section in slides_source:
        slide_idx += 1
        slide = prs.slides.add_slide(bullet_slide_layout)
        try:
            slide.shapes.title.text = section.get("heading", section.get("title", "")) 
            _apply_formatting(slide.shapes.title, is_title=True)
        except Exception:
            pass 
        
        # Heuristic to find body placeholder
        tf = None
        for shape in slide.placeholders:
             if shape.placeholder_format.idx == 1:
                 tf = shape.text_frame
                 break
        
        if tf:
            tf.clear() 
            for bullet in section["bullets"]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                # Apply style immediately?
                p.font.name = target_font
                p.font.size = Pt(18)
                p.font.color.rgb = _hex_to_rgb(body_color_hex)
        
        _add_slide_number(slide, slide_idx)
            
    # Patent Claims Slide
    slide_idx += 1
    slide = prs.slides.add_slide(bullet_slide_layout)
    try:
        slide.shapes.title.text = "Patent-style Claims"
        _apply_formatting(slide.shapes.title, is_title=True)
    except: pass
    
    tf = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            break
            
    if tf:
        tf.clear()
        for claim in plan["claims"]:
            p = tf.add_paragraph()
            p.text = claim
            p.level = 0
            p.font.name = target_font
            p.font.size = Pt(16)
            p.font.color.rgb = _hex_to_rgb(body_color_hex)
            
    _add_slide_number(slide, slide_idx)

    # Post-processing: Backgrounds if requested (Simplified)
    if bg_image_path:
        for slide in prs.slides:
             try:
                 # Hacky, Z-order issues persist in python-pptx without XML manipulation.
                 # We accept it might overlay or be overlaid.
                 pic = slide.shapes.add_picture(str(bg_image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
                 slide.shapes._spTree.remove(pic._element)
                 slide.shapes._spTree.insert(0, pic._element)
             except: pass
            
    prs.save(str(path))


# ... inside chat, update signature ...

# ... inside chat, update signature ...

@app.post("/regenerate")
async def regenerate(req: RegenerateRequest):
    plan = req.plan
    settings = req.ppt_settings
    cert_settings = req.cert_settings

    # We can't easily preserve the original query message for patent/certificate generation 
    # if it's not in the plan. For now, we'll use a placeholder or extract from summary.
    message = plan.get("summary", "2025") # Used as Date/Context placeholder 

    report_path = FILES_DIR / "report.pdf"
    report_docx_path = FILES_DIR / "report.docx"
    ppt_path = FILES_DIR / "slides.pptx"
    patent_path = FILES_DIR / "patent.pdf"
    cert_zip_path = FILES_DIR / "certificates.zip"

    try:
        # Defaults for regenerate
        docx_tmpl_path = None
        pptx_tmpl_path = None
        
        _build_report_pdf(plan, report_path)
        try:
            _build_docx(plan, report_docx_path, template_path=docx_tmpl_path)
        except Exception as e: 
            print(f"DOCX Build Error: {e}")
            traceback.print_exc()

        _build_ppt(plan, ppt_path, settings=settings, template_path=pptx_tmpl_path) 
        
        _build_patent_pdf(plan, message, patent_path)
        _build_certificates_zip(plan, message, cert_zip_path, settings=cert_settings)
        
        return {
            "status": "ok", 
            "files": _file_urls(), 
            "plan": plan, 
            "reply": "Project regenerated.", 
            "ai": {"enabled": True, "note": "Regenerated from plan"}
        }
    except Exception as e:
        traceback.print_exc()
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/ai_assist")
async def ai_assist(req: AiAssistRequest):
    eff_provider = req.provider or DEFAULT_AI_PROVIDER
    eff_key = req.api_key or os.getenv("OPENAI_API_KEY") 
    eff_model = req.model or DEFAULT_AI_MODEL
    
    if not (eff_provider and eff_key):
        return JSONResponse({"error": "No AI credentials provided"}, status_code=400)

    client = _get_start_client(eff_provider, eff_key, eff_model)
    if not client:
        return JSONResponse({"error": "Failed to init AI client"}, status_code=500)

    prompt = f"Instruction: {req.instruction}\nText: {req.text}\nOutput:"
    
    try:
        # Simplified call - mostly reusing logic or direct call
        # OpenAI/OpenRouter
        if eff_provider in {"openai", "openrouter"}:
            resp = client.chat.completions.create(
                model=eff_model,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=500
            )
            return {"text": resp.choices[0].message.content.strip()}
            
        # Gemini
        if eff_provider == "gemini":
            model = client.GenerativeModel(eff_model)
            resp = model.generate_content(prompt)
            return {"text": resp.text.strip()}
            
        # Others not implemented for assist yet
        return {"text": req.text} 
        
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)



def _file_urls():
    """Generates the file download URLs."""
    ts = int(time.time())
    return {
        "report": f"/files/report.pdf?t={ts}",
        "report_docx": f"/files/report.docx?t={ts}",
        "ppt": f"/files/slides.pptx?t={ts}",
        "patent": f"/files/patent.pdf?t={ts}",
        "certificates": f"/files/certificates.zip?t={ts}",
    }

@app.post("/chat")
async def chat(
    message: str = Form(...),
    file: UploadFile | None = None,
    template_file: UploadFile | None = None, # For Certificates
    template_docx: UploadFile | None = None, # For Word Reports
    template_pptx: UploadFile | None = None, # For PPT Slides
    mode: str = Form("static"),
    api_key: str = Form(None),
    provider: str = Form(None),
    model: str = Form(None),
    settings_json: str = Form(None),
    cert_settings: str = Form(None),
    language: str = Form("en"),
):
    # Parse cert_settings if provided
    cert_opts = {}
    if cert_settings:
        try:
            cert_opts = json.loads(cert_settings)
        except: pass
    
    # Handle CSV
    csv_data = []
    if file and file.filename.lower().endswith(".csv"):
        try:
            content = await file.read()
            if content:
                text = content.decode("utf-8")
                f_io = io.StringIO(text)
                # Check if has header
                sample = text[:1024]
                try:
                    has_header = csv.Sniffer().has_header(sample)
                except:
                    has_header = False
                
                if has_header:
                    f_io.seek(0)
                    reader = csv.DictReader(f_io)
                    csv_data = [row for row in reader]
                else:
                    # Fallback to simple split logic if no header (treat col 0 as name)
                    lines = text.strip().splitlines()
                    for line in lines:
                        parts = line.split(",")
                        if parts and parts[0].strip():
                            csv_data.append({"name": parts[0].strip()})
        except Exception as e:
            print(f"CSV Parsing Error: {e}")
            traceback.print_exc()

    # Handle Templates
    docx_tmpl_path = None
    if template_docx:
        docx_tmpl_path = FILES_DIR / f"user_template_{int(time.time())}.docx"
        with open(docx_tmpl_path, "wb") as f:
            f.write(await template_docx.read())

    pptx_tmpl_path = None
    if template_pptx:
        pptx_tmpl_path = FILES_DIR / f"user_template_{int(time.time())}.pptx"
        with open(pptx_tmpl_path, "wb") as f:
            f.write(await template_pptx.read())

    # custom cert template (image)
    if template_file:
         try:
            tmpl_content = await template_file.read()
            custom_tmpl_path = FILES_DIR / f"custom_template_{template_file.filename}"
            with open(custom_tmpl_path, "wb") as f:
                f.write(tmpl_content)
            if cert_opts:
                cert_opts["background_image"] = base64.b64encode(tmpl_content).decode("utf-8")
         except Exception as e:
            print(f"Template upload failed: {e}")

    ai_error: str | None = None
    ai_used = False
    is_greet = _is_greeting(message)
    
    if is_greet:
        plan = _greeting_plan(message)
    else:
        plan = _fallback_plan(message, file)

    if mode == "hybrid" and not is_greet:
        eff_provider = provider or DEFAULT_AI_PROVIDER
        eff_key = api_key or os.getenv("OPENAI_API_KEY") 
        eff_model = model or DEFAULT_AI_MODEL
        
        # GitHub Context Injection
        gh_url = _extract_github_url(message)
        if gh_url:
            gh_context = _fetch_github_context(gh_url)
            message = f"{message}\n\n{gh_context}"

        # Language Instruction
        if language and language != "en":
             message = f"{message}\n\nOUTPUT RULE: Translate all user-facing content (summary, bullets, claims) to {language}. Keep JSON keys (title, sections, heading, bullets, claims) in English."

        if eff_provider and eff_key:
             ai_plan, ai_error = _call_ai_plan_dynamic(message, eff_provider, eff_key, eff_model)
             if ai_plan:
                 plan = _normalize_plan(ai_plan, message)
                 # Add note about GitHub context
                 if gh_url:
                     plan["summary"] += f" (Analyzed GitHub Repo: {gh_url})"
                 ai_used = True
        else:
             ai_error = "Missing API Key or Provider configuration"

    file_urls = None
    should_gen_cert = False
    if not is_greet:
        report_path = FILES_DIR / "report.pdf"
        report_docx_path = FILES_DIR / "report.docx"
        ppt_path = FILES_DIR / "slides.pptx"
        patent_path = FILES_DIR / "patent.pdf"
        cert_zip_path = FILES_DIR / "certificates.zip"

        _build_report_pdf(plan, report_path)
        try:
            _build_docx(plan, report_docx_path, template_path=docx_tmpl_path)
        except Exception as e:
            print(f"DOCX Build Error: {e}")
            
        _build_ppt(plan, ppt_path, template_path=pptx_tmpl_path)
        _build_patent_pdf(plan, message, patent_path)
        
        # Conditional Certificate Generation
        msg_lower = message.lower()
        should_gen_cert = bool(csv_data) or any(k in msg_lower for k in ["certificate", "cert ", "diploma"])
        
        if should_gen_cert:
            _build_certificates_zip(plan, message, cert_zip_path, data=csv_data, settings=cert_opts)
        
        file_urls = _file_urls()
        if not should_gen_cert:
             file_urls.pop("certificates", None)

    reply_lines = [plan["summary"]]
    if ai_error:
        reply_lines.append(f"AI warning: {ai_error}")
    
    if csv_data:
        reply_lines.append(f"\nGenerated {len(csv_data)} certificates from CSV.")
    elif should_gen_cert and not is_greet:
         reply_lines.append("\nCertificates generated.")

    return JSONResponse(
        {
            "reply": "\n".join(reply_lines),
            "files": file_urls,
            "plan": plan,
            "ai": {
                "enabled": ai_used,
                "provider": provider,
                "model": model,
                "mode_used": "hybrid" if ai_used else "static",
                "error": ai_error,
            },
        }
    )


@app.get("/files/{filename}")
async def get_file(filename: str):
    file_path = FILES_DIR / filename
    if file_path.exists():
        return FileResponse(file_path)
    return JSONResponse({"error": "File not found"}, status_code=404)


@app.get("/status")
async def status():
    # Deprecated mostly for dynamic usage, but kept for handshake
    return {"message": "Server Ready"}
