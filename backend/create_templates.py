from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_modern_template():
    prs = Presentation()
    
    # 1. Slide Master modifications (simulated by creating a custom layout on the fly usually, 
    # but here we just create a file we can load later)
    # Actually python-pptx doesn't support creating Masters easily. 
    # We will create a presentation with pre-defined slide styles on the first slides 
    # and use `prs.slides.add_slide()` repeatedly? No, we need a true template.
    # Workaround: We will use this script to CREATE a 'template_modern.pptx' 
    # which has a dark background and specific font styles pre-set on the master layouts if possible.
    # Since python-pptx is limited in *editing* masters, we will create a "base" pptx 
    # that we will validly open and "save as" new files.
    
    # Set background color for the master (hacky, by iterating layouts)
    # This is complex in python-pptx. 
    # Simpler approach: Create a PPTX with one Title Slide and one Content Slide
    # configured with the look we want, and use that as our "Base".

    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # We can't easily change the master background via python-pptx without deep XML hacking.
    # So we will implement a "design application" step in the main app instead.
    # BUT, we can save this file as a placeholder.
    
    prs.save("backend/templates/template_modern.pptx")
    print("Created empty template_modern.pptx")

if __name__ == "__main__":
    create_modern_template()
